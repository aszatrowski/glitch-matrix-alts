from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime

prs = Presentation()

# title slide
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Phenotype Covariance Ongoing Slides"
for paragraph in title.text_frame.paragraphs:
    paragraph.font.name = "Helvetica Neue"
    paragraph.font.size = Pt(44)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0, 51, 102)  # dark blue

subtitle.text = "Austin Szatrowski | " + str(datetime.today().strftime('%Y-%m-%d'))
for paragraph in subtitle.text_frame.paragraphs:
    paragraph.font.name = "Helvetica Neue"
    paragraph.font.size = Pt(30)
    paragraph.font.color.rgb = RGBColor(179, 189, 199)  # dark blue


blank_slide = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_slide)
top = Inches(1.5)
width = Inches(8)
left = Inches(1)

txBox = slide2.shapes.add_textbox(left = Inches(0.5), top = Inches(0.25), width = width, height = Inches(1))
tf = txBox.text_frame
for paragraph in tf.paragraphs:
    paragraph.text = "Curves don't separate at 5 generations"
    paragraph.font.name = "Helvetica Neue"
    paragraph.font.size = Pt(20)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0, 51, 102)  # dark blue
slide2.shapes.add_picture(snakemake.input[0], left, top, width=width)

slide3 = prs.slides.add_slide(blank_slide)
left = Inches(1)
top = Inches(2)
slide3.shapes.add_picture(snakemake.input[1], left, top, width=width)

slide4 = prs.slides.add_slide(blank_slide)
left = Inches(1)
top = Inches(2)
slide4.shapes.add_picture(snakemake.input[2], left, top, width=width)

prs.save(snakemake.output['pptx'])